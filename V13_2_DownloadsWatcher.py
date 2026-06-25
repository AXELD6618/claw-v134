#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
V13.2 浏览器下载文件夹 — 知识管道扩展数据源 (Downloads Watcher)
============================================================
举一反三：将Chrome/Edge下载文件夹作为知识管道的附加数据源。
扫描新下载的研报、投资分析、行业文档，自动提取内容和分类，接入KnowledgePipeline。

方案优势:
  - 零用户操作 — 下载即自动接入知识管道
  - 复用ContentClassifier — 与WeChatFileSync共享分类引擎
  - 去重机制 — 基于文件哈希，不重复处理
  - 多来源标记 — source='browser_download'，区分于wechat_file

版本: V13.2 DW-1
创建: 2026-06-24
作者: 毕方灵犀·天眼 (亚瑟)
"""

import os
import sys
import json
import sqlite3
import time
import hashlib
import logging
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Optional, Set

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, BASE_DIR)

os.makedirs('logs', exist_ok=True)
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('logs/downloads_watcher.log', encoding='utf-8'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger('DownloadsWatcher')

# Reuse the ContentClassifier from WeChatFileSync
from V13_2_WeChatFileSync import SyncConfig, ContentClassifier


class DownloadsWatcher:
    """浏览器下载文件夹监听器 — 扩展知识管道数据源"""
    
    # 监控的文件类型（研报/文档类）
    WATCH_EXTENSIONS = {'.pdf', '.docx', '.doc', '.pptx', '.ppt', '.xlsx', '.xls',
                         '.txt', '.html', '.htm', '.md', '.csv'}
    
    # 排除模式（非研报/非投资类文件）
    EXCLUDE_PATTERNS = ['installer', 'setup', '.exe', '.msi', '.zip', '.rar', '.7z',
                        '.vsix', '.skill', 'desktop.ini', '.lnk', '.mp4', '.mov',
                        '.mp3', '.jpg', '.jpeg', '.png', '.gif']
    
    def __init__(self, download_dirs: List[str] = None, db_path: str = None):
        self.config = SyncConfig()
        self.classifier = ContentClassifier(self.config)
        self.db_path = db_path or self.config.db_path
        
        # 监控目录
        if download_dirs is None:
            home = os.path.expanduser('~')
            self.download_dirs = [
                os.path.join(home, 'Downloads'),
            ]
        else:
            self.download_dirs = download_dirs
        
        # 过滤存在的目录
        self.download_dirs = [d for d in self.download_dirs if os.path.isdir(d)]
        
        self._processed_hashes: Set[str] = set()
        self._init_db()
        self._load_processed_hashes()
        
        logger.info(f"[DownloadsWatcher] 初始化完成")
        logger.info(f"[DownloadsWatcher] 监控目录: {self.download_dirs}")
        logger.info(f"[DownloadsWatcher] 已处理文件: {len(self._processed_hashes)}")
    
    # ============================================================
    # 数据库
    # ============================================================
    
    def _init_db(self):
        """初始化下载记录表"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            cursor.execute("""
                CREATE TABLE IF NOT EXISTS downloads_sync (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    file_path TEXT NOT NULL,
                    file_name TEXT NOT NULL,
                    file_size INTEGER,
                    file_hash TEXT UNIQUE NOT NULL,
                    file_type TEXT,
                    content_extracted TEXT,
                    content_hash TEXT,
                    extracted_chars INTEGER DEFAULT 0,
                    category TEXT DEFAULT 'other',
                    category_confidence REAL DEFAULT 0,
                    keywords_matched TEXT,
                    source_dir TEXT,
                    synced_at TEXT DEFAULT CURRENT_TIMESTAMP,
                    knowledge_pipeline_fed INTEGER DEFAULT 0
                )
            """)
            cursor.execute("""
                CREATE INDEX IF NOT EXISTS idx_dl_hash ON downloads_sync(file_hash)
            """)
            cursor.execute("""
                CREATE INDEX IF NOT EXISTS idx_dl_synced ON downloads_sync(synced_at)
            """)
            conn.commit()
            conn.close()
            logger.info("[DownloadsWatcher] 数据库表已初始化")
        except Exception as e:
            logger.error(f"[DownloadsWatcher] 数据库初始化失败: {e}")
    
    def _load_processed_hashes(self):
        """加载已处理的文件哈希"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            cursor.execute("SELECT file_hash FROM downloads_sync")
            self._processed_hashes = {row[0] for row in cursor.fetchall()}
            conn.close()
        except Exception as e:
            logger.error(f"[DownloadsWatcher] 加载哈希失败: {e}")
            self._processed_hashes = set()
    
    # ============================================================
    # 文件哈希
    # ============================================================
    
    def _compute_file_hash(self, file_path: str) -> str:
        """计算文件哈希"""
        file_size = os.path.getsize(file_path)
        try:
            if file_size < 5 * 1024 * 1024:
                sha256 = hashlib.sha256()
                with open(file_path, 'rb') as f:
                    for chunk in iter(lambda: f.read(65536), b''):
                        sha256.update(chunk)
                return sha256.hexdigest()
            else:
                mtime = os.path.getmtime(file_path)
                name = os.path.basename(file_path)
                md5 = hashlib.md5()
                with open(file_path, 'rb') as f:
                    md5.update(f.read(65536))
                md5.update(f"{name}{file_size}{mtime}".encode())
                return md5.hexdigest()
        except Exception as e:
            logger.error(f"[DownloadsWatcher] 哈希失败 {file_path}: {e}")
            return ""
    
    # ============================================================
    # 内容提取
    # ============================================================
    
    def _extract_text_file(self, file_path: str, max_chars: int = 50000) -> Optional[str]:
        """文本文件提取"""
        encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
        for enc in encodings:
            try:
                with open(file_path, 'r', encoding=enc, errors='replace') as f:
                    return f.read(max_chars)
            except:
                continue
        return None
    
    def _extract_pdf(self, file_path: str, max_chars: int = 50000) -> Optional[str]:
        """PDF提取"""
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            texts = []
            total = 0
            for page in reader.pages:
                pt = page.extract_text()
                if pt:
                    texts.append(pt)
                    total += len(pt)
                    if total >= max_chars:
                        break
            content = '\n'.join(texts)
            return content[:max_chars] if content.strip() else None
        except:
            pass
        
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                texts = []
                total = 0
                for page in pdf.pages:
                    pt = page.extract_text()
                    if pt:
                        texts.append(pt)
                        total += len(pt)
                        if total >= max_chars:
                            break
                content = '\n'.join(texts)
                return content[:max_chars] if content.strip() else None
        except:
            pass
        return None
    
    def _extract_docx(self, file_path: str, max_chars: int = 50000) -> Optional[str]:
        """DOCX提取"""
        try:
            from docx import Document
            doc = Document(file_path)
            texts = []
            total = 0
            for para in doc.paragraphs:
                if para.text.strip():
                    texts.append(para.text)
                    total += len(para.text)
                    if total >= max_chars:
                        break
            content = '\n'.join(texts)
            return content[:max_chars] if content.strip() else None
        except:
            return None
    
    def _extract_pptx(self, file_path: str, max_chars: int = 50000) -> Optional[str]:
        """PPT提取"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            texts = []
            total = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        texts.append(shape.text)
                        total += len(shape.text)
                        if total >= max_chars:
                            break
            content = '\n'.join(texts)
            return content[:max_chars] if content.strip() else None
        except:
            return None
    
    def _extract_xlsx(self, file_path: str, max_chars: int = 50000) -> Optional[str]:
        """Excel提取"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            texts = []
            total = 0
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                texts.append(f"[Sheet: {sheet_name}]")
                for row in ws.iter_rows(values_only=True):
                    row_text = ' | '.join([str(c) if c is not None else '' for c in row])
                    if row_text.strip():
                        texts.append(row_text)
                        total += len(row_text)
                        if total >= max_chars:
                            break
                if total >= max_chars:
                    break
            wb.close()
            return '\n'.join(texts)[:max_chars] if texts else f"[{file_path} - 元数据: sheets={len(wb.sheetnames)}]"
        except:
            return None
    
    def _extract_html(self, file_path: str, max_chars: int = 50000) -> Optional[str]:
        """HTML提取（纯文本）"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                html = f.read(max_chars * 2)
            from html.parser import HTMLParser
            class TextExtractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.texts = []
                    self.skip = {'script', 'style', 'meta', 'link'}
                    self.current_tag = None
                def handle_starttag(self, tag, attrs):
                    self.current_tag = tag
                def handle_endtag(self, tag):
                    self.current_tag = None
                def handle_data(self, data):
                    if self.current_tag not in self.skip:
                        t = data.strip()
                        if t:
                            self.texts.append(t)
            extractor = TextExtractor()
            extractor.feed(html)
            return '\n'.join(extractor.texts)[:max_chars]
        except:
            return self._extract_text_file(file_path, max_chars)
    
    def _extract_content(self, file_path: str, ext: str) -> Optional[str]:
        """根据扩展名提取内容"""
        ext = ext.lower()
        if ext in ['.txt', '.md', '.csv', '.log', '.json', '.xml']:
            return self._extract_text_file(file_path)
        elif ext in ['.pdf']:
            return self._extract_pdf(file_path)
        elif ext in ['.docx', '.doc']:
            return self._extract_docx(file_path)
        elif ext in ['.pptx', '.ppt']:
            return self._extract_pptx(file_path)
        elif ext in ['.xlsx', '.xls']:
            return self._extract_xlsx(file_path)
        elif ext in ['.html', '.htm']:
            return self._extract_html(file_path)
        return None
    
    # ============================================================
    # 扫描与处理
    # ============================================================
    
    def _should_process(self, file_path: str) -> bool:
        """判断是否应处理此文件"""
        fname = os.path.basename(file_path).lower()
        ext = os.path.splitext(file_path)[1].lower()
        
        # 扩展名过滤
        if ext not in self.WATCH_EXTENSIONS:
            return False
        
        # 文件名排除模式
        for pattern in self.EXCLUDE_PATTERNS:
            if pattern.lower() in fname:
                return False
        
        # 大小检查
        try:
            if os.path.getsize(file_path) > 100 * 1024 * 1024:  # 100MB max
                return False
        except:
            return False
        
        return True
    
    def scan(self) -> Dict:
        """扫描下载文件夹，处理新文件"""
        stats = {'scanned': 0, 'new': 0, 'skipped': 0, 'error': 0,
                 'classified': {'knowledge': 0, 'trading': 0, 'sentiment': 0, 'other': 0}}
        new_items = []
        
        for dl_dir in self.download_dirs:
            if not os.path.isdir(dl_dir):
                continue
            
            for root, dirs, files in os.walk(dl_dir):
                # 只扫描一层（不递归子目录）
                if root != dl_dir:
                    dirs.clear()
                    continue
                
                for fname in files:
                    file_path = os.path.join(root, fname)
                    stats['scanned'] += 1
                    
                    if not self._should_process(file_path):
                        stats['skipped'] += 1
                        continue
                    
                    # 哈希去重
                    file_hash = self._compute_file_hash(file_path)
                    if not file_hash or file_hash in self._processed_hashes:
                        stats['skipped'] += 1
                        continue
                    
                    # 提取内容
                    ext = os.path.splitext(fname)[1]
                    size = os.path.getsize(file_path)
                    
                    try:
                        content = self._extract_content(file_path, ext)
                        content_hash = hashlib.md5((content or '')[:2000].encode()).hexdigest()
                        
                        # 分类
                        result = self.classifier.classify(fname, content or '', ext)
                        
                        # 存入数据库
                        conn = sqlite3.connect(self.db_path)
                        cur = conn.cursor()
                        cur.execute("""
                            INSERT INTO downloads_sync 
                            (file_path, file_name, file_size, file_hash, file_type,
                             content_extracted, content_hash, extracted_chars,
                             category, category_confidence, keywords_matched, source_dir)
                            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                        """, (
                            file_path, fname, size, file_hash, ext,
                            content or '', content_hash, len(content or ''),
                            result['category'], result['confidence'],
                            json.dumps(result.get('keywords_matched', []), ensure_ascii=False),
                            dl_dir
                        ))
                        conn.commit()
                        conn.close()
                        
                        self._processed_hashes.add(file_hash)
                        stats['new'] += 1
                        stats['classified'][result['category']] += 1
                        
                        if result['category'] in ('knowledge', 'trading'):
                            new_items.append({
                                'file_path': file_path,
                                'file_name': fname,
                                'content': content or '',
                                'category': result['category'],
                                'confidence': result['confidence'],
                                'keywords': result.get('keywords_matched', []),
                            })
                        
                        logger.info(f"[DownloadsWatcher] 新文件: {fname[:60]} → {result['label']}({result['confidence']:.0%})")
                        
                    except Exception as e:
                        logger.error(f"[DownloadsWatcher] 处理失败 {fname}: {e}")
                        stats['error'] += 1
        
        # 接入知识管道
        if new_items:
            try:
                self._feed_knowledge_pipeline(new_items)
            except Exception as e:
                logger.error(f"[DownloadsWatcher] 知识管道接入失败: {e}")
        
        return stats
    
    def _feed_knowledge_pipeline(self, items: List[Dict]) -> int:
        """将知识类下载文件接入KnowledgePipeline"""
        try:
            from V13_2_KnowledgePipeline import KnowledgePipeline
            kp = KnowledgePipeline()
            count = kp.ingest_wechat_file([{
                'source': 'browser_download',
                'type': item['file_type'] if 'file_type' in item else 'file',
                'title': item.get('file_name', ''),
                'content': item.get('content', ''),
                'published_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            } for item in items])
            
            # 更新数据库标记
            conn = sqlite3.connect(self.db_path)
            cur = conn.cursor()
            cur.execute("UPDATE downloads_sync SET knowledge_pipeline_fed=1 WHERE knowledge_pipeline_fed=0")
            conn.commit()
            conn.close()
            
            logger.info(f"[DownloadsWatcher] 知识管道接入: {count}条")
            return count
        except Exception as e:
            logger.warning(f"[DownloadsWatcher] 知识管道接入失败(非关键): {e}")
            return 0
    
    def get_stats(self) -> Dict:
        """获取统计信息"""
        try:
            conn = sqlite3.connect(self.db_path)
            cur = conn.cursor()
            cur.execute("SELECT COUNT(*) FROM downloads_sync")
            total = cur.fetchone()[0]
            cur.execute("SELECT category, COUNT(*) FROM downloads_sync GROUP BY category")
            categories = {row[0]: row[1] for row in cur.fetchall()}
            cur.execute("SELECT COUNT(*) FROM downloads_sync WHERE knowledge_pipeline_fed=1")
            fed = cur.fetchone()[0]
            conn.close()
            return {
                'total': total,
                'fed_to_pipeline': fed,
                'categories': categories,
            }
        except:
            return {'total': 0, 'fed_to_pipeline': 0, 'categories': {}}


# ============================================================
# 独立运行入口
# ============================================================

def main():
    """独立运行：扫描下载文件夹"""
    logger.info("=" * 60)
    logger.info("DownloadsWatcher 启动扫描")
    
    watcher = DownloadsWatcher()
    
    # 扫描
    stats = watcher.scan()
    logger.info(f"扫描完成: {stats}")
    
    # 统计
    db_stats = watcher.get_stats()
    logger.info(f"数据库统计: {db_stats}")
    
    logger.info("=" * 60)
    return stats


if __name__ == '__main__':
    main()
